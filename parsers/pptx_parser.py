from __future__ import annotations

import io
import os

from pptx import Presentation
from pptx.util import Emu
from PIL import Image

from .base import ParsedDocument, ParsedPage
from utils.logger import get_logger

logger = get_logger(__name__)


def parse_pptx(file_path: str) -> ParsedDocument:
    """Extract text, tables, and image OCR from a PPTX file."""
    prs = Presentation(file_path)
    pages: list[ParsedPage] = []

    for slide_idx, slide in enumerate(prs.slides):
        text_parts: list[str] = []
        tables: list[list[list[str]]] = []
        images_text: list[str] = []

        for shape in slide.shapes:
            # Text frames (titles, body text, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text_parts.append(line)

            # Tables
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                tables.append(rows)

            # Images
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    image_blob = shape.image.blob
                    img = Image.open(io.BytesIO(image_blob))
                    w, h = img.size
                    if w >= 100 and h >= 100:
                        ocr_text = _ocr_image(img)
                        if ocr_text.strip():
                            images_text.append(ocr_text)
                except Exception as e:
                    logger.warning(
                        f"PPTX image extraction failed on slide {slide_idx + 1}: {e}"
                    )

        # Get slide notes
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            text_parts.append(f"[演讲者备注] {notes}")

        pages.append(
            ParsedPage(
                page_number=slide_idx + 1,
                text="\n".join(text_parts),
                tables=tables,
                images_text=images_text,
            )
        )

    logger.info(f"PPTX parsed: {file_path} ({len(pages)} slides)")
    return ParsedDocument(
        file_name=os.path.basename(file_path),
        file_path=file_path,
        file_type="pptx",
        pages=pages,
        metadata={"total_slides": len(pages)},
    )


def _ocr_image(img: Image.Image) -> str:
    try:
        import pytesseract

        return pytesseract.image_to_string(img, lang="chi_sim+eng")
    except Exception:
        return ""
